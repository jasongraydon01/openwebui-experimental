import os
import time
import hashlib
import re
import tempfile
import subprocess
import json
from collections import Counter
from typing import List, Dict, Any

# External dependencies
import pinecone
from pinecone import ServerlessSpec
import pptx
import requests
import sqlite3
from dotenv import load_dotenv
# New dependencies
from docling.document_converter import DocumentConverter

# -------------------------------
# Configuration and Initialization
# -------------------------------

# Load environment variables
load_dotenv()

# Environment variables
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
DATABASE_PATH = os.getenv("DATABASE_PATH")
VLLM_EMBED_URL = os.getenv("VLLM_EMBED_URL")
VLLM_CHAT_URL = os.getenv("VLLM_CHAT_URL")
FOLDER_PATH = os.getenv("FOLDER_PATH")

# Validate critical environment variables
if not PINECONE_API_KEY:
    raise ValueError("Pinecone API key is missing. Please set PINECONE_API_KEY in your .env file.")
if not PINECONE_INDEX_NAME:
    raise ValueError("Pinecone index name is missing. Please set PINECONE_INDEX_NAME in your .env file.")
if not DATABASE_PATH:
    raise ValueError("Database path is missing. Please set DATABASE_PATH in your .env file.")
if not FOLDER_PATH:
    raise ValueError("Folder path is missing. Please set FOLDER_PATH in your .env file.")
if not VLLM_EMBED_URL:
    raise ValueError("vLLM embed URL is missing. Please set VLLM_EMBED_URL in your .env file.")
if not VLLM_CHAT_URL:
    raise ValueError("vLLM chat URL is missing. Please set VLLM_CHAT_URL in your .env file.")

# Initialize Pinecone client
pc = pinecone.Pinecone(api_key=PINECONE_API_KEY)

# Create index if it doesn't exist
if not pc.has_index(PINECONE_INDEX_NAME):
    pc.create_index(
        name=PINECONE_INDEX_NAME,
        dimension=3584,
        metric="cosine",
        spec=ServerlessSpec(cloud="aws", region="us-east-1")
    )
    print(f"Created new Pinecone index: {PINECONE_INDEX_NAME}")

# Wait for index to be ready
while not pc.describe_index(PINECONE_INDEX_NAME).status['ready']:
    print("Waiting for Pinecone index to be ready...")
    time.sleep(1)

# Connect to the index
index = pc.Index(PINECONE_INDEX_NAME)
print(f"Connected to Pinecone index: {PINECONE_INDEX_NAME}")

# -------------------------------
# Database Functions
# -------------------------------

def init_database():
    """Initialize the SQLite database with necessary tables."""
    conn = sqlite3.connect(DATABASE_PATH)
    c = conn.cursor()
    
    # Create file_log table if it doesn't exist
    c.execute('''
    CREATE TABLE IF NOT EXISTS file_log (
        file_name TEXT PRIMARY KEY,
        file_hash TEXT,
        last_modified REAL,
        last_processed REAL,
        slide_count INTEGER,
        research_type TEXT,
        project_type TEXT,
        client TEXT,
        product TEXT
    )
    ''')
    
    # Create slides table if it doesn't exist
    c.execute('''
    CREATE TABLE IF NOT EXISTS slides (
        vector_id TEXT PRIMARY KEY,
        file_name TEXT,
        slide_number INTEGER,
        content TEXT,
        summary TEXT,
        keywords TEXT,
        raw_json TEXT,
        FOREIGN KEY (file_name) REFERENCES file_log (file_name) ON DELETE CASCADE
    )
    ''')
    
    # Check for and add new columns if needed
    columns_to_add = {
        "file_log": ["research_type", "project_type", "client", "product"],
        "slides": ["summary", "raw_json"]
    }
    
    for table, new_columns in columns_to_add.items():
        c.execute(f"PRAGMA table_info({table})")
        existing_columns = [column[1] for column in c.fetchall()]
        
        for column in new_columns:
            if column not in existing_columns:
                c.execute(f"ALTER TABLE {table} ADD COLUMN {column} TEXT")
                print(f"Added {column} column to {table} table")
    
    conn.commit()
    conn.close()
    print("Database initialized successfully")

def remove_file_data(file_name):
    """Remove all data associated with a file from the database and Pinecone."""
    conn = sqlite3.connect(DATABASE_PATH)
    c = conn.cursor()
    
    # Get all vector IDs for this file
    c.execute("SELECT vector_id FROM slides WHERE file_name=?", (file_name,))
    vector_ids = [row[0] for row in c.fetchall()]
    
    # Remove the file record and slides from the database
    c.execute("DELETE FROM file_log WHERE file_name=?", (file_name,))
    c.execute("DELETE FROM slides WHERE file_name=?", (file_name,))
    conn.commit()
    conn.close()
    
    # Remove all vectors associated with this file from Pinecone
    if vector_ids:
        # Process in batches of 100 to avoid overwhelming the API
        for i in range(0, len(vector_ids), 100):
            batch = vector_ids[i:i+100]
            index.delete(ids=batch, namespace="ns1")
        print(f"Removed {len(vector_ids)} vectors for {file_name} from Pinecone")
    
    return len(vector_ids) if vector_ids else 0

def check_for_removed_files():
    """Check for files that have been removed and clean up their data."""
    current_files = set(os.listdir(FOLDER_PATH))
    
    conn = sqlite3.connect(DATABASE_PATH)
    c = conn.cursor()
    c.execute("SELECT file_name FROM file_log")
    files_in_db = set(row[0] for row in c.fetchall())
    conn.close()
    
    # Determine which files have been removed
    removed_files = files_in_db - current_files
    
    if removed_files:
        print(f"Detected {len(removed_files)} removed files")
        for file in removed_files:
            vectors_removed = remove_file_data(file)
            print(f"Cleaned up {file}: removed {vectors_removed} vectors")
    
    return removed_files

# -------------------------------
# New LLM Functions
# -------------------------------

def categorize_presentation(presentation_content):
    """
    Categorize a presentation based on its content using an LLM.
    
    Args:
        presentation_content: A string containing the content of the entire presentation
        
    Returns:
        A dictionary with categorization information
    """
    # Prepare the prompt for the LLM
    prompt = f"""
    Analyze the following presentation content and categorize it across these dimensions:
    - Research Type: Qualitative or Quantitative
    - Project Type: Segmentation, ATU, Demand Study, Message Testing, etc.
    - Client: Identify the pharmaceutical client (e.g., Pfizer, JNJ, etc.)
    - Product: Identify the product (e.g., Spravato, Fintepla, etc.)
    
    Return your analysis in JSON format with these four keys: "research_type", "project_type", "client", "product". 
    
    Presentation Content:
    {presentation_content[:20000]}  # Limiting to first 20000 chars to avoid token limits
    """
    
    try:
        # Use the vLLM API to get the categorization
        payload = {
            "prompt": prompt,
            "model": "Qwen/Qwen2.5-7B-Instruct",
            "temperature": 0.2,  # Low temperature for more consistent results
            "max_tokens": 500
        }
        response = requests.post(VLLM_CHAT_URL, json=payload, timeout=60)
        response.raise_for_status()
        
        # Parse the response to extract the JSON
        llm_response = response.json().get("text", "") 
        
        # Extract JSON from the response
        try:
            json_str = llm_response.strip()
            if json_str.startswith('```json'):
                json_str = json_str.replace('```json', '').replace('```', '')
            elif json_str.startswith('```'):
                json_str = json_str.replace('```', '')
            categorization = json.loads(json_str.strip())
        except json.JSONDecodeError:
            # Fallback: try to extract JSON-like structure from text
            import re
            json_pattern = r'{\s*"[^"]+"\s*:\s*"[^"]+"\s*,\s*"[^"]+"\s*:\s*"[^"]+"\s*,\s*"[^"]+"\s*:\s*"[^"]+"\s*,\s*"[^"]+"\s*:\s*"[^"]+"\s*}'
            match = re.search(json_pattern, llm_response)
            if match:
                categorization = json.loads(match.group(0))
            else:
                # If all else fails, create a default categorization
                categorization = {
                    "research_type": "Unknown",
                    "project_type": "Unknown",
                    "client": "Unknown",
                    "product": "Unknown"
                }
        
        # Ensure all expected keys are present
        expected_keys = ["research_type", "project_type", "client", "product"]
        for key in expected_keys:
            if key not in categorization:
                categorization[key] = "Unknown"
        
        return categorization
    
    except Exception as e:
        print(f"Error categorizing presentation: {e}")
        # Return default categorization on error
        return {
            "research_type": "Unknown",
            "project_type": "Unknown",
            "client": "Unknown",
            "product": "Unknown"
        }

def summarize_slide(slide_content):
    """
    Generate a concise summary of a slide using an LLM.
    
    Args:
        slide_content: A dictionary containing the slide's content
        
    Returns:
        A string containing the summary
    """
    # Prepare the content for summarization
    title = slide_content.get("title", "")
    text = slide_content.get("text", "")
    
    # Convert tables to text representation
    table_text = ""
    for table in slide_content.get("tables", []):
        for row in table:
            table_text += " | ".join(row) + "\n"
    
    # Combine all content
    full_content = f"Title: {title}\n\nContent: {text}\n\nTables: {table_text}"
    
    # Prepare the system message and user message for the chat API
    system_message = """You summarize presentation slides concisely while preserving all numerical data exactly as presented. Follow these guidelines:
    1. Provide a clear overview of the main points and insights
    2. For any numerical data, statistics or percentages, quote them EXACTLY as they appear in the original text
    3. Preserve all specific metrics, numbers and quantitative findings in their original form"""
    
    user_message = f"Summarize this slide content:\n\n{full_content}"
    
    try:
        # Use the vLLM API to get the summary using the chat completions format
        payload = {
            "model": "Qwen/Qwen2.5-7B-Instruct",
            "messages": [
                {"role": "system", "content": system_message},
                {"role": "user", "content": user_message}
            ],
            "temperature": 0.2  # Lower temperature for more precise data extraction
        }
        response = requests.post(VLLM_CHAT_URL, json=payload, timeout=30)
        response.raise_for_status()
        
        # Extract the summary from the response - chat completions use a different response format
        response_data = response.json()
        summary = response_data.get("choices", [{}])[0].get("message", {}).get("content", "").strip()
        
        return summary
    
    except Exception as e:
        print(f"Error summarizing slide: {e}")
        # Return original text as fallback
        return f"Title: {title}\n\n{text[:500]}..."

# -------------------------------
# New File Processing Functions
# -------------------------------

def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Convert PPTX to PDF using libreoffice."""
    libreoffice_path = '/usr/lib/libreoffice/program/libmergedlo.so'
    subprocess.run([libreoffice_path, '--headless', '--convert-to', 'pdf', pptx_path, '--outdir', os.path.dirname(pdf_path)])
    print(f'Converted {pptx_path} to {pdf_path}')

def extract_pptx_content_enhanced(pptx_path):
    """
    Extract content from PowerPoint file using pptxtopdf and docling for improved accuracy.
    
    Args:
        pptx_path: Path to the PowerPoint file
        
    Returns:
        A list of dictionaries, each containing information about a slide
    """
    slides_data = []
    
    try:
        # Step 1: Convert PPTX to PDF using pptxtopdf
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as pdf_file:
            pdf_path = pdf_file.name
        
        # Convert PPTX to PDF
        convert_pptx_to_pdf(pptx_path, pdf_path)
        print(f'Converted {pptx_path} to {pdf_path}')
        
        # Step 2: Use DocumentConverter to process the PDF
        converter = DocumentConverter()
        result = converter.convert(pdf_path)
        
        # Step 3: Extract content from the document
        pages = result.document.pages
        
        for i, page in enumerate(pages):
            slide_number = i + 1
            
            # Extract title (first heading or first text block)
            slide_title = ""
            for block in page.blocks:
                if block.type == "heading" or (not slide_title and block.type == "text"):
                    slide_title = block.text
                    break
            
            # Extract all text content
            slide_text = "\n".join([block.text for block in page.blocks if block.type == "text" or block.type == "heading"])
            
            # Extract tables
            tables = []
            for block in page.blocks:
                if block.type == "table":
                    table_data = []
                    for row in block.rows:
                        table_row = [cell.text for cell in row.cells]
                        table_data.append(table_row)
                    tables.append(table_data)
            
            # Store the raw page data as JSON
            raw_json = json.dumps(page.to_dict() if hasattr(page, "to_dict") else {})
            
            slides_data.append({
                "slide_number": slide_number,
                "title": slide_title or f"Slide {slide_number}",
                "text": slide_text,
                "tables": tables,
                "raw_json": raw_json
            })
        
        # Clean up temporary file
        os.unlink(pdf_path)
        
    except Exception as e:
        print(f"Error in enhanced extraction for {pptx_path}: {e}")
        # Fall back to the original extraction method
        print(f"Falling back to original extraction method for {pptx_path}")
        slides_data = extract_pptx_content(pptx_path)
        # Add raw_json field with empty value
        for slide in slides_data:
            slide["raw_json"] = "{}"
    
    return slides_data
# -------------------------------
# Original File Processing Functions (Updated)
# -------------------------------

def get_file_hash(file_path):
    """Calculate MD5 hash of a file for change detection."""
    hasher = hashlib.md5()
    with open(file_path, 'rb') as f:
        buf = f.read(65536)
        while len(buf) > 0:
            hasher.update(buf)
            buf = f.read(65536)
    return hasher.hexdigest()

def extract_pptx_content(pptx_path):
    """Extract content from PowerPoint file using python-pptx."""
    prs = pptx.Presentation(pptx_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        # Extract slide title
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text
        
        # Extract text from all shapes
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text += shape.text + "\n"
        
        # Extract tables
        tables = []
        for shape in slide.shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
                tables.append(table_data)

        slides_data.append({
            "slide_number": i + 1,
            "title": slide_title,
            "text": slide_text,
            "tables": tables,
            "raw_json": "{}"  # Empty raw_json for compatibility
        })
    
    return slides_data

def extract_keywords(text, tables, max_keywords=15):
    """Extract keywords from text using a statistical approach similar to TF-IDF."""
    # Combine all text
    table_text = ""
    for table in tables:
        table_text += " ".join([" ".join(row) for row in table])
    
    all_text = f"{text} {table_text}"
    
    # Clean and normalize the text
    all_text = all_text.lower()
    all_text = re.sub(r'[^\w\s]', ' ', all_text)
    
    # Tokenize
    words = all_text.split()
    
    # Common stopwords to filter out
    stopwords = {'a', 'an', 'the', 'and', 'or', 'but', 'if', 'because', 'as', 'what', 
                'which', 'this', 'that', 'these', 'those', 'then', 'just', 'so', 'than', 'such',
                'when', 'who', 'how', 'where', 'why', 'is', 'are', 'was', 'were', 'be', 'been',
                'have', 'has', 'had', 'do', 'does', 'did', 'but', 'at', 'by', 'with', 'from',
                'here', 'there', 'to', 'of', 'for', 'in', 'on', 'about', 'into', 'over', 'after'}
    
    # Filter words
    filtered_words = [word for word in words if word not in stopwords and len(word) > 2]
    
    # Get word frequencies
    word_counts = Counter(filtered_words)
    
    # Get top single keywords
    keywords = [word for word, count in word_counts.most_common(max_keywords)]
    
    # Extract important bigrams (2-word phrases)
    bigrams = []
    for i in range(len(words) - 1):
        if (words[i] not in stopwords and words[i+1] not in stopwords and 
            len(words[i]) > 2 and len(words[i+1]) > 2):
            bigrams.append(f"{words[i]} {words[i+1]}")
    
    bigram_counts = Counter(bigrams)
    top_bigrams = [bigram for bigram, count in bigram_counts.most_common(5)]
    
    # Combine keywords and bigrams, limiting to max_keywords total
    all_keywords = keywords[:10] + top_bigrams
    
    return ", ".join(all_keywords[:max_keywords])

def generate_embedding(text):
    """Generate embedding vector using vLLM API."""
    try:
        payload = {
            "input": text,
            "model": "Alibaba-NLP/gte-Qwen2-7B-instruct"
        }
        response = requests.post(VLLM_EMBED_URL, json=payload, timeout=30)
        response.raise_for_status()
        return response.json()["data"][0]["embedding"]
    except requests.exceptions.RequestException as e:
        print(f"Error generating embedding: {e}")
        raise

def process_slide(slide, file_name):
    """Process a single slide without context.
    
    Args:
        slide: The slide data to process
        file_name: Name of the PowerPoint file
    """
    slide_number = slide["slide_number"]
    slide_text = slide["text"]
    slide_title = slide["title"] or f"Slide {slide_number}"
    table_content = slide["tables"]
    
    # Convert slide's table content to text
    table_text = ""
    for table in table_content:
        table_text += " ".join([" ".join(row) for row in table])
    
    # Build content for the slide
    full_content = f"Title: {slide_title}\n\nContent: {slide_text}"
    if table_text:
        full_content += f"\n\nTables: {table_text}"
    
    # Extract keywords from slide
    keywords = extract_keywords(slide_text, table_content)
    
    # Create vector ID
    vector_id = f"{file_name}_{slide_number}"
    
    return {
        "vector_id": vector_id,
        "slide_number": slide_number,
        "content": full_content,
        "keywords": keywords,
        "raw_json": slide.get("raw_json", "{}")  # Include raw JSON
    }

# -------------------------------
# Updated Main Processing Function
# -------------------------------

def process_pptx_files():
    """Main function to process PowerPoint files and update the vector database."""
    # Initialize database and check for removed files
    init_database()
    check_for_removed_files()
    
    # Track vectors to upload
    vectors_to_upload = []
    
    # Connect to database
    conn = sqlite3.connect(DATABASE_PATH)
    c = conn.cursor()
    
    # Process each PowerPoint file in the folder
    for pptx_file in os.listdir(FOLDER_PATH):
        if not pptx_file.endswith(".pptx"):
            continue
            
        pptx_path = os.path.join(FOLDER_PATH, pptx_file)
        file_hash = get_file_hash(pptx_path)
        last_modified = os.path.getmtime(pptx_path)
        
        # Check if file has changed
        c.execute("SELECT file_hash FROM file_log WHERE file_name=?", (pptx_file,))
        result = c.fetchone()
        
        if result and result[0] == file_hash:
            print(f"Skipping {pptx_file}: No changes detected")
            continue
            
        # If file has changed or is new, process it
        if result:
            # Remove existing data first
            vectors_removed = remove_file_data(pptx_file)
            print(f"Updating {pptx_file}: Removed {vectors_removed} previous vectors")
        else:
            print(f"Processing new file: {pptx_file}")
            
        # Extract content from PowerPoint using the enhanced method
        slides_data = extract_pptx_content_enhanced(pptx_path)
        print(f"Extracted {len(slides_data)} slides from {pptx_file}")
        
        # Categorize the presentation
        # Combine all slide content for categorization
        presentation_content = "\n\n".join([
            f"Slide {slide['slide_number']}: {slide['title']}\n{slide['text']}"
            for slide in slides_data
        ])
        
        categorization = categorize_presentation(presentation_content)
        print(f"Categorized {pptx_file} as: {categorization}")
        
        # Process each slide individually
        for slide in slides_data:
            # Generate summary for the current slide
            slide_summary = summarize_slide(slide)
            
            # Process slide without context
            processed_slide = process_slide(slide, pptx_file)
            
            # Add summary to the processed slide
            processed_slide["summary"] = slide_summary
            
            # Generate embedding for the summary rather than the full content
            embedding = generate_embedding(slide_summary)
            
            # Prepare metadata - ensure all values are of proper types for Pinecone
            metadata = {
                "file_name": pptx_file,
                "slide_number": processed_slide["slide_number"],
                "one_drive_link": f"https://onedrive.com/{pptx_file}",
                "content_preview": slide_summary[:1000],  # Use summary for preview
                "keywords": processed_slide["keywords"],
                "last_modified": float(last_modified),
                "research_type": categorization["research_type"],
                "project_type": categorization["project_type"],
                "client": categorization["client"],
                "product": categorization["product"]
            }
            
            # Add to vectors batch
            vectors_to_upload.append({
                "id": processed_slide["vector_id"],
                "values": embedding,
                "metadata": metadata
            })
            
            # Store in database
            c.execute(
                "INSERT INTO slides (vector_id, file_name, slide_number, content, summary, keywords, raw_json) VALUES (?, ?, ?, ?, ?, ?, ?)",
                (
                    processed_slide["vector_id"],
                    pptx_file,
                    processed_slide["slide_number"],
                    processed_slide["content"],
                    slide_summary,
                    processed_slide["keywords"],
                    processed_slide["raw_json"]
                )
            )
        
        # Update file log with categorization
        c.execute(
            "REPLACE INTO file_log (file_name, file_hash, last_modified, last_processed, slide_count, research_type, project_type, client, product) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)",
            (
                pptx_file, 
                file_hash, 
                last_modified, 
                time.time(), 
                len(slides_data),
                categorization["research_type"],
                categorization["project_type"],
                categorization["client"],
                categorization["product"]
            )
        )
        
        print(f"Processed {pptx_file}: {len(slides_data)} slides")
    
    # Commit database changes
    conn.commit()
    conn.close()
    
    # Upload vectors to Pinecone in batches
    if vectors_to_upload:
        total_vectors = len(vectors_to_upload)
        batch_size = 50  # Optimal batch size for Pinecone
        
        for i in range(0, total_vectors, batch_size):
            batch = vectors_to_upload[i:i+batch_size]
            index.upsert(vectors=batch, namespace="ns1")
            print(f"Uploaded batch {i//batch_size + 1}/{(total_vectors-1)//batch_size + 1}: {len(batch)} vectors")
        
        print(f"Successfully uploaded {total_vectors} vectors to Pinecone")
    else:
        print("No new or updated files to process")

# -------------------------------
# Main Execution
# -------------------------------

if __name__ == "__main__":
    print("Starting PowerPoint processing pipeline")
    process_pptx_files()
    print("Processing complete")